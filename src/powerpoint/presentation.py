import os
import re

from pptx import Presentation

from src.powerpoint.slide_enum import SlideEnum
from src.utils.common_scripts import clean_up_string


# PowerPointPresentation class
# This class is used to create a PowerPoint presentation
# It uses the python-pptx library to create a presentation
# The class has methods to add slides, save the presentation and get slides
# The class also has methods to get a slide layout, and get a slide by slide number


class PowerPointPresentation:
    """ PowerPointPresentation class to handle all PowerPoint Presentation API calls
    """

    def __init__(self, presentation_string=None, presentation_theme=None):
        """ Constructor for PowerPointPresentation class
        :param presentation_string: The presentation string
        :param presentation_theme: The presentation theme
        """
        if presentation_theme == "Blank" or presentation_theme == "":
            self.theme = None
        else:
            self.theme = presentation_theme

        cwd = os.getcwd()

        if not self.theme:
            self.presentation = Presentation()
        elif not os.path.exists(cwd + "/static/PresentationThemes/"):
            self.presentation = Presentation()
        elif self.theme in os.listdir(cwd + "/static/PresentationThemes/"):
            self.presentation = Presentation(pptx=cwd + "/static/PresentationThemes/" + self.theme)
        elif self.theme not in (cwd + "/static/PresentationThemes/"):
            self.presentation = Presentation()

        if presentation_string:
            self.populate_presentation(presentation_string)

    # Return a string representation of the class
    def __str__(self):
        """ Returns the string representation of the PowerPointPresentation class
        :return: The string representation of the PowerPointPresentation class
        """
        return "PowerPointPresentation"

    @classmethod
    def _slice_presentation(cls, presliced_presentation):
        """ Private Method - returns a sliced presentation
        :param presliced_presentation: The presentation to slice
        :return: The sliced presentation
        :raises ValueError: If the presentation is empty
        """

        if not presliced_presentation:
            raise ValueError("Presentation cannot be empty.")
        else:
            sliced_presentation = []

            for presentation_line in presliced_presentation.splitlines():

                presentation_line = clean_up_string(presentation_line)

                if presentation_line.lower().startswith("slide"):
                    # If the presentation line starts with Slide, then it is the start of a new slide
                    # Add the line to the sliced presentation as a new slide, including a new line, so it can be
                    # formatted correctly. Done via regex to replace "Slide x" with "Title:"
                    slide_title = re.sub(r"Slide \d:", "TITLE:", presentation_line, flags=re.IGNORECASE)
                    sliced_presentation.append(slide_title + "\n")
                elif len(sliced_presentation) != 0:
                    # Otherwise, it is part of the previous slide
                    # Add the line to the previous slide, including a new line, so it can be formatted correctly
                    sliced_presentation[-1] += presentation_line + "\n"
                else:
                    # Otherwise, it doesn't start with Slide, but is the first line of the presentation
                    # This may happen on the first slide of the presentation
                    # So set it as the first line in the sliced presentation
                    sliced_presentation.append(presentation_line + "\n")

            return sliced_presentation

    @classmethod
    def _get_slide_content(cls, slide_pages):
        """ Private Method - returns a title, subtitle, content and notes for each slide - if they exist
        :param slide_pages: The slide pages
        :return: A title, subtitle, content and notes for each slide
        :raises ValueError: If the slide pages are empty
        """

        if not slide_pages:
            raise ValueError("Slide pages cannot be empty.")
        else:
            slide_title = ""
            slide_subtitle = ""
            slide_content = ""
            slide_notes = ""
            slide_image = ""
            last_section = ""
            for section in slide_pages.splitlines():
                section = clean_up_string(section)

                if len(section) == 0:
                    continue

                # parse the reply into the different sections
                if section.lower().startswith("TITLE".casefold()):
                    slide_title = section[6:].strip()
                elif section.lower().startswith("SUBTITLE".casefold()):
                    slide_subtitle = section[9:].strip()
                elif section.lower().startswith("CONTENT".casefold()):
                    # If the slide content is "Content:", then set the slide content to the section, without the
                    # "Content:" prefix and the "-" suffix, and strip any whitespace
                    # This is done to ensure that the slide content is not set to "Content:"
                    slide_content = section[8:].replace("-", "").strip() + "\n"
                    # If the slide content is "(No content ", then set the slide content to an empty string
                    # This is done to ensure that the slide content is not set to "(No content required)"
                    if "(no content" in slide_content.lower():
                        slide_content = ""
                    last_section = "Content"
                elif section.lower().startswith("NOTES".casefold()):
                    slide_notes = section[6:].replace("-", "").strip() + "\n"
                    last_section = "Notes"
                elif section.lower().startswith("IMAGE".casefold()):
                    slide_image = section[6:].strip()
                elif section.lower().startswith("REFERENCES".casefold()):
                    # If the section starts with "References", insert this section into the slide content
                    # This is done to ensure that the slide content is not set to "References"
                    # removes the references header and adds it to the content
                    slide_content = section[10:].strip() + "\n"
                    last_section = "Content"
                elif section.lower().startswith("SCRIPT".casefold()):
                    # If the section starts with "Script", insert this section into the slide notes
                    # This is done incase the response doesn't have a notes section, or it is separated from the
                    # slide content
                    slide_notes += section.strip() + "\n"
                    last_section = "Notes"
                elif last_section == "Content":
                    slide_content += section.replace("-", "").strip() + "\n"
                elif last_section == "Notes":
                    slide_notes += section.replace("-", "").strip() + "\n"

            return slide_title, slide_subtitle, slide_content, slide_notes, slide_image

    @classmethod
    def _set_layouts(cls, title, subtitle, content, image, notes):
        """ Private Method - returns the best slide layout for the presentation
        :param title: The title of the slide
        :param subtitle: The subtitle of the slide
        :param content: The content of the slide
        :param image: The image of the slide
        :param notes: The notes of the slide
        :return: The best slide layout for the presentation
        """

        if title and subtitle:
            if image:
                # if there is an image, then use the picture with caption layout
                return SlideEnum.SLIDE_PICTURE_WITH_CAPTION_LAYOUT.value
            return SlideEnum.SLIDE_TITLE_LAYOUT.value
        elif image:
            return SlideEnum.SLIDE_PICTURE_WITH_CAPTION_LAYOUT.value
        elif title and content:
            return SlideEnum.SLIDE_TITLE_AND_CONTENT_LAYOUT.value
        elif title:
            return SlideEnum.SLIDE_TITLE_ONLY_LAYOUT.value
        elif content:
            return SlideEnum.SLIDE_CONTENT_WITH_CAPTION_LAYOUT.value
        else:
            return SlideEnum.SLIDE_BLANK_LAYOUT.value

    def populate_presentation(self, input_presentation):
        """ Returns a populated presentation
        :param input_presentation: The input presentation
        :return: The populated presentation
        """

        sliced_presentation = self._slice_presentation(input_presentation)
        for index, section in enumerate(sliced_presentation):

            if len(section) == 0:
                continue
            title, subtitle, content, notes, image = self._get_slide_content(section)

            if index == 0:
                # for the first slide, we don't want any images; as we want to show the title instead
                image = None

            if index == 0 and self.theme:
                first_slide_with_theme = True
            else:
                first_slide_with_theme = False

            self.add_slide(title=title, subtitle=subtitle, text=content, notes=notes, image=image,
                           first_slide_with_theme=first_slide_with_theme)

    def add_slide(self, title=None, subtitle=None, text=None, image=None, notes=None, first_slide_with_theme=False):
        """ Returns the created slide
        :param title: The title of the slide
        :param subtitle: The subtitle of the slide
        :param text: The text of the slide
        :param image: The image of the slide
        :param notes: The notes of the slide
        :param first_slide_with_theme: The first slide with the theme, meaning we want to use this slide rather than add
        a new one
        :return: The created slide
        """

        slide_layout = self._set_layouts(title, subtitle, text, image, notes)

        created_slide = self.create_slide(slide_layout, first_slide_with_theme)

        if slide_layout == SlideEnum.SLIDE_PICTURE_WITH_CAPTION_LAYOUT.value:
            created_slide = self._create_slides_with_images(created_slide, title, subtitle, text, image, notes)
        else:
            created_slide = self._create_slides_without_images(created_slide, title, subtitle, text, notes)
        return created_slide

    @classmethod
    def _create_slides_with_images(cls, created_slide, title, subtitle, text, image, notes):
        """ Returns the created slide with images
        :param created_slide: The created slide
        :param title: The title of the slide
        :param subtitle: The subtitle of the slide
        :param text: The text of the slide
        :param image: The image of the slide
        :param notes: The notes of the slide
        :return: The created slide with images
        """

        if image:
            created_slide.placeholders[1].insert_picture(image)
        else:
            raise FileNotFoundError("Image must be provided for picture with caption layout.")
        if title:
            created_slide.shapes.title.text = title
        if subtitle:
            created_slide.placeholders[1].text = subtitle
        if text:
            created_slide.placeholders[2].text = text
        if notes:
            created_slide.notes_slide.notes_text_frame.text = notes
        return created_slide

    @classmethod
    def _create_slides_without_images(cls, created_slide, title, subtitle, text, notes):
        """ Returns the created slide without images
        :param created_slide: The created slide
        :param title: The title of the slide
        :param subtitle: The subtitle of the slide
        :param text: The text of the slide
        :param notes: The notes of the slide
        :return: The created slide without images
        """

        if title:
            created_slide.shapes.title.text = title
        if subtitle:
            created_slide.placeholders[1].text = subtitle
        if text:
            created_slide.placeholders[1].text = text
        if notes:
            created_slide.notes_slide.notes_text_frame.text = notes
        return created_slide

    def save(self, filename):
        """ Returns the saved presentation
        :param filename: The filename of the presentation
        :return: The saved presentation
        :raises ValueError: If the filename is empty or does not end with .pptx
        """

        if not filename:
            raise ValueError("Filename cannot be empty.")
        else:
            return self.presentation.save(filename)

    def get_slide_layout(self, slide_layout):
        """ Returns the slide layout
        :param slide_layout: The slide layout
        :return: The slide layout
        :raises ValueError: If the slide layout is not between 0 and 8 or not an integer
        """

        if int(slide_layout) < 0 or int(slide_layout) > (len(SlideEnum) - 1):
            raise ValueError("Slide layout must be between 0 and " + str(len(SlideEnum) - 1) + ".")
        elif not isinstance(slide_layout, int):
            raise ValueError("Slide layout must be an integer.")
        else:
            return self.presentation.slide_layouts[slide_layout]

    def get_slide(self, slide_number):
        """ Returns the slide specified by the slide number
        :param slide_number: The slide number
        :return: The slide
        :raises ValueError: If the slide number is not between 0 and the number of slides in the presentation or not
        an integer
        """

        if not isinstance(slide_number, int):
            raise ValueError("Slide number must be an integer.")
        elif slide_number < 0 or slide_number > len(self.presentation.slides):
            raise ValueError("Slide number must be between 0 and the number of slides in the presentation.")
        else:
            return self.presentation.slides[slide_number]

    def get_slides(self):
        """ Returns all slides in the presentation
        :return: All slides in the presentation
        """

        return self.presentation.slides

    def create_slide(self, slide_layout, first_slide_with_theme=False):
        """ Returns the created slide
        :param slide_layout: The slide layout
        :param first_slide_with_theme: The first slide with the theme, meaning we want to use this slide rather than add
        a new one
        :return: The created slide
        """

        if first_slide_with_theme:
            return self.get_slide(0)
        else:
            return self.presentation.slides.add_slide(self.presentation.slide_layouts[slide_layout])
